#%%
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor

class UVMNode:
    def __init__(self, name, child, type) -> None:
        self.name = name
        self.child = child
        self.type = type

class MakeUVMPowerPoint:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = 16 * 914400  
        self.prs.slide_height = 9 * 914400
        self.slide = None
        self.slide_layout = None

    def CreateSlide(self):
        # Add a slide
        self.slide_layout = self.prs.slide_layouts[0]  # Select a slide layout (0 is usually the title slide layout)
        self.slide = self.prs.slides.add_slide(self.slide_layout)

    def CreateComponent(self, x, y, width, height, name, 
                        shape = MSO_SHAPE.ROUNDED_RECTANGLE,
                        P_alignment = PP_ALIGN.CENTER,
                        color = RGBColor(0, 160, 0)):
        left   = self.prs.slide_width  * x
        top    = self.prs.slide_height * y
        width  = self.prs.slide_width  * width
        height = self.prs.slide_height * height
        oval   = self.slide.shapes.add_shape(shape, left, top, width, height)
        
        text_box = oval.text_frame
        text_box.text = name
        text_box.paragraphs[0].alignment = P_alignment
        text_box.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

        fill = oval.fill
        fill.solid()
        fill.fore_color.rgb = color
        return None

    def SavePowerPoint(self):
        # Save the presentation
        self.prs.save("my_presentation.pptx")

    def CreateTest(self, x = 0.1, y = 0.1, NumberAgent = 1, name = "Test"):
        self.CreateComponent(x = x, y = y, width = 0.075 + 0.225 * NumberAgent, height = 0.5, name = name,
                             shape = MSO_SHAPE.RECTANGLE,
                             P_alignment = PP_ALIGN.LEFT,
                             color = RGBColor(112, 48, 160))
        return None
    def CreateEnv(self, x = 0.1, y = 0.1, NumberAgent = 1, name = "Env"):
        self.CreateComponent(x = x + 0.025, y = y + 0.05, width = 0.025 + 0.225 * NumberAgent, height = 0.4, name = name,
                             shape = MSO_SHAPE.RECTANGLE,
                             P_alignment = PP_ALIGN.LEFT,
                             color = RGBColor(0,112,192))
        return None
    
    def CreateAgent(self, x = 0.1, y = 0.1, NumberAgent = 1, name = "Agent"):
        for i in range(NumberAgent):
            self.CreateComponent(x = x + 0.05 + i * 0.225, y = y + 0.2, width = 0.2, height = 0.2, name = name,
                                shape = MSO_SHAPE.RECTANGLE,
                                P_alignment = PP_ALIGN.LEFT,
                                color = RGBColor(0, 176, 240))
            self.CreateComponent(x = x + 0.075 + i * 0.225, y = y + 0.3, width = 0.07, height = 0.05, name = "Driver")
            self.CreateComponent(x = x + 0.165 + i * 0.225, y = y + 0.3, width = 0.07, height = 0.05, name = "Monitor")
    
    def CalculateAgent(self, test_root):
        NumberAgent = 0
        #BFS
        queue = [test_root]

        while queue:
            temp = queue.pop(0)
            if temp.type == "agent":
                NumberAgent = NumberAgent + 1
            if temp.child == None:
                continue
            for c in temp.child:
                queue.append(c)
        return NumberAgent
    
    def RunByBFS(self, test_root):
        NumberAgent = self.CalculateAgent(test_root)
        print(NumberAgent)
        x = 0.5 - (0.075 + 0.225 * NumberAgent)/2
        y = 0.2
        self.CreateSlide()
        queue = [test_root]

        while queue:
            temp = queue.pop(0)
            if temp.type == "test":
                self.CreateTest(x, y, NumberAgent = NumberAgent, name = temp.name)
            if temp.type == "env":
                self.CreateEnv(x,y, NumberAgent = NumberAgent, name = temp.name)
            if temp.type == "scoreboard":
                self.CreateComponent(x =  0.45, y = y + 0.1, width = 0.1, height = 0.05, name = temp.name,
                             color = RGBColor(0, 176, 240))
            if temp.type == "agent":
                self.CreateAgent(x,y,NumberAgent = NumberAgent, name = temp.name)
            if temp.child == None:
                continue
            for c in temp.child:
                queue.append(c)
        self.SavePowerPoint()
        return None
    
    def Run(self):
        self.CreateSlide()
        NumberAgent = 2
        x = 0.5 - (0.075 + 0.225 * NumberAgent)/2
        y = 0.2
        self.CreateTest(x, y, NumberAgent = NumberAgent)
        self.CreateEnv(x,y, NumberAgent = NumberAgent)
        self.CreateComponent(x =  0.45, y = y + 0.1, width = 0.1, height = 0.05, name = "Scoreboard",
                             color = RGBColor(0, 176, 240))
        self.CreateAgent(x,y,NumberAgent = NumberAgent)

        self.SavePowerPoint()
"""
if __name__ == "__main__":
    
    MUP = MakeUVMPowerPoint()
    MUP.Run()
"""    
#%%





